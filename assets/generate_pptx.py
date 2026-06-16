"""Generate AI Recruiter Agent Swarm pitch deck as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

C = {
    'bg': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x1E, 0x29, 0x3B),
    'muted': RGBColor(0x64, 0x74, 0x8B),
    'light': RGBColor(0x94, 0xA3, 0xB8),
    'accent': RGBColor(0x4F, 0x46, 0xE5),
    'green': RGBColor(0x16, 0xA3, 0x4A),
    'card_bg': RGBColor(0xF8, 0xFA, 0xFC),
    'card_border': RGBColor(0xE2, 0xE8, 0xF0),
    'amber': RGBColor(0xD9, 0x77, 0x06),
    'pink': RGBColor(0xDB, 0x27, 0x77),
    'cyan': RGBColor(0x08, 0x91, 0xB2),
    'red': RGBColor(0xDC, 0x26, 0x26),
}

def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C['bg']
    return slide

def rect(slide, left, top, width, height, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border:
        shape.line.fill.solid()
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def txtbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=14, bold=False, color=C['text'], align=PP_ALIGN.LEFT, font_name='Calibri'):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return p

def add_para(tf, text, size=14, bold=False, color=C['text'], align=PP_ALIGN.LEFT, space_before=0, font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    return p

# ---------- SLIDE 1: Title ----------
s = add_slide(prs)

# Tag boxes
rect(s, Inches(0.6), Inches(0.6), Inches(1.8), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.62), Inches(1.6), Inches(0.3))
set_text(tb.text_frame, 'AI SYSTEMS ARCHITECT', 10, bold=True, color=C['accent'])

rect(s, Inches(2.6), Inches(0.6), Inches(1.4), Inches(0.32), fill=RGBColor(0xF0, 0xFD, 0xF4))
tb = txtbox(s, Inches(2.7), Inches(0.62), Inches(1.2), Inches(0.3))
set_text(tb.text_frame, 'INDIA RUNS 2026', 10, bold=True, color=C['green'])

# Title
tb = txtbox(s, Inches(0.6), Inches(1.4), Inches(8), Inches(1.6))
tf = tb.text_frame
set_text(tf, 'AI Recruiter', 48, bold=True, color=C['text'])
add_para(tf, 'Agent Swarm', 48, bold=True, color=C['accent'])

# Subtitle
tb = txtbox(s, Inches(0.6), Inches(3.2), Inches(8), Inches(0.5))
set_text(tb.text_frame, 'Autonomous multi-agent hiring for the next India', 20, color=C['muted'])

# Footer
tb = txtbox(s, Inches(0.6), Inches(5.6), Inches(4), Inches(0.8))
tf = tb.text_frame
set_text(tf, 'peter14l', 14, bold=True, color=C['text'])
add_para(tf, 'India Runs by Redrob AI', 12, color=C['light'])

# Number
tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '1/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 2: Problem ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'THE PROBLEM', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.7))
tf = tb.text_frame
set_text(tf, "India's hiring is ", 32, bold=True, color=C['text'])
add_para(tf, 'broken at scale', 32, bold=True, color=C['accent'])

# 3 stat cards
stat_data = [
    ('500+', 'Applications per role', "95% irrelevant — recruiters drown in noise"),
    ('60%', 'Time spent screening', "Not interviewing. Not hiring. Just filtering."),
    ('Hidden\u2193', 'Non-obvious talent', "Tier 2/3, non-IIT, non-English candidates get missed"),
]
for i, (num, label, desc) in enumerate(stat_data):
    x = Inches(0.6 + i * 3.8)
    y = Inches(2.0)
    rect(s, x, y, Inches(3.5), Inches(2.0), fill=C['card_bg'], border=C['card_border'])
    tb = txtbox(s, x + Inches(0.2), y + Inches(0.15), Inches(3.1), Inches(1.7))
    tf = tb.text_frame
    set_text(tf, num, 30, bold=True, color=C['accent'] if i != 2 else C['green'])
    add_para(tf, label, 11, bold=True, color=C['light'])
    add_para(tf, desc, 13, color=C['muted'], space_before=6)

# Warning box
y2 = Inches(4.2)
rect(s, Inches(0.6), y2, Inches(11.5), Inches(1.0), fill=C['card_bg'], border=RGBColor(0xFE, 0xCA, 0xCA))
tb = txtbox(s, Inches(0.9), y2 + Inches(0.1), Inches(11), Inches(0.8))
tf = tb.text_frame
set_text(tf, 'Current AI tools are black boxes', 14, bold=True, color=C['red'])
add_para(tf, 'They reject without explanation. No transparency, no trust, no feedback loop. Keyword filters miss transferable skills entirely.', 12, color=C['muted'])

tb = txtbox(s, Inches(0.6), Inches(5.5), Inches(8), Inches(0.4))
set_text(tb.text_frame, "India's recruitment market: $2B+ and growing 15% YoY — core tech hasn't changed in a decade.", 11, color=C['light'])

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '2/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 3: Why Now ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.0), Inches(0.32), fill=RGBColor(0xF0, 0xFD, 0xF4))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(0.8), Inches(0.3))
set_text(tb.text_frame, 'WHY NOW', 10, bold=True, color=C['green'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.7))
tf = tb.text_frame
set_text(tf, 'Three forces make this ', 32, bold=True, color=C['text'])
add_para(tf, 'possible today', 32, bold=True, color=C['green'])

forces = [
    ('AI Maturity', "LLMs reason, plan, and use tools. Multi-agent frameworks are production-ready (CrewAI, AutoGen). Voice AI in Indian languages works (Sarvam, Bhashini)."),
    ("India's Talent Boom", "1.5M+ engineers graduate annually. Tier 2/3 cities produce world-class talent. Remote work demands location-agnostic hiring."),
    ("Redrob's Mission", "Building India's AI OS for work. AI-native hiring is the foundational layer. Built for India, in India — not a Silicon Valley port."),
]
for i, (title, desc) in enumerate(forces):
    x = Inches(0.6 + i * 4.0)
    y = Inches(2.0)
    rect(s, x, y, Inches(3.7), Inches(2.5), fill=C['card_bg'], border=C['card_border'])
    tb = txtbox(s, x + Inches(0.2), y + Inches(0.15), Inches(3.3), Inches(2.2))
    tf = tb.text_frame
    set_text(tf, title, 16, bold=True, color=C['text'])
    add_para(tf, desc, 13, color=C['muted'], space_before=8)

# Bottom bar
y3 = Inches(4.8)
rect(s, Inches(0.6), y3, Inches(11.5), Inches(0.6), fill=C['card_bg'], border=C['card_border'])
tb = txtbox(s, Inches(0.9), y3 + Inches(0.1), Inches(11), Inches(0.4))
set_text(tb.text_frame, 'This combination of AI maturity, market need, and platform timing has never existed before.', 13, bold=True, color=C['muted'])

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '3/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 4: Solution ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'THE SOLUTION', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.6))
tf = tb.text_frame
set_text(tf, 'Enter the ', 32, bold=True, color=C['text'])
add_para(tf, 'Agent Swarm', 32, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.7), Inches(8), Inches(0.4))
set_text(tb.text_frame, '5 specialized AI agents + 1 Orchestrator. From JD to shortlist in under 10 minutes.', 14, color=C['muted'])

# Agent chips in table
agents = [
    ('Orchestrator', 'Brain', 'Routes tasks, manages pipeline state, decides when to escalate to humans', C['accent']),
    ('Sourcer', 'Scout', 'Searches LinkedIn, GitHub, Naukri, internal databases autonomously', C['green']),
    ('Matcher', 'Judge', 'Semantic ranking, transferable skill detection, explainable scoring', C['amber']),
    ('Screener', 'Interviewer', 'Voice calls in Hinglish, dynamic questioning, sentiment analysis', C['pink']),
    ('Scheduler', 'Coordinator', 'Calendar sync, slot negotiation, reminders, feedback collection', C['cyan']),
]

# Table header
x0, y0 = Inches(0.6), Inches(2.3)
w_cols = [Inches(1.8), Inches(0.9), Inches(7.5)]
rect(s, x0, y0, w_cols[0]+w_cols[1]+w_cols[2], Inches(0.35), fill=C['card_bg'], border=C['card_border'])
tb = txtbox(s, x0 + Inches(0.15), y0 + Inches(0.05), w_cols[0], Inches(0.3))
set_text(tb.text_frame, 'Agent', 10, bold=True, color=C['light'])
tb = txtbox(s, x0 + w_cols[0] + Inches(0.15), y0 + Inches(0.05), w_cols[1], Inches(0.3))
set_text(tb.text_frame, 'Role', 10, bold=True, color=C['light'])
tb = txtbox(s, x0 + w_cols[0] + w_cols[1] + Inches(0.15), y0 + Inches(0.05), w_cols[2], Inches(0.3))
set_text(tb.text_frame, 'Responsibility', 10, bold=True, color=C['light'])

for i, (name, role, desc, clr) in enumerate(agents):
    y = y0 + Inches(0.35) + i * Inches(0.38)
    rect(s, x0, y, w_cols[0]+w_cols[1]+w_cols[2], Inches(0.38), border=C['card_border'])
    tb = txtbox(s, x0 + Inches(0.15), y + Inches(0.06), w_cols[0], Inches(0.3))
    set_text(tb.text_frame, name, 13, bold=True, color=clr)
    tb = txtbox(s, x0 + w_cols[0] + Inches(0.15), y + Inches(0.06), w_cols[1], Inches(0.3))
    set_text(tb.text_frame, role, 12, color=C['muted'])
    tb = txtbox(s, x0 + w_cols[0] + w_cols[1] + Inches(0.15), y + Inches(0.06), w_cols[2], Inches(0.3))
    set_text(tb.text_frame, desc, 12, color=C['light'])

# Footer emphasis
tb = txtbox(s, Inches(0.6), Inches(5.8), Inches(8), Inches(0.4))
set_text(tb.text_frame, 'Every decision explained. Every candidate gets a fair shot. No black boxes.', 12, bold=True, color=C['accent'])

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '4/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 5: How It Works ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'HOW IT WORKS', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.7))
set_text(tb.text_frame, 'From JD to shortlist in under 10 minutes', 30, bold=True, color=C['text'])

steps_left = [
    ('1', 'JD Submitted', 'Recruiter pastes job description. Orchestrator parses and activates the swarm.'),
    ('2', 'Multi-Source Sourcing', 'Sourcer Agent searches 4+ sources simultaneously — LinkedIn, GitHub, Naukri, internal DB.'),
    ('3', 'Semantic Matching', 'Matcher Agent scores candidates across 5 dimensions with full transparency.'),
]
steps_right = [
    ('4', 'Voice Screening', 'Screener Agent conducts natural voice calls in Hinglish — adaptive, contextual, fair.'),
    ('5', 'Auto-Scheduling', 'Scheduler Agent books interviews in one click. No email tennis required.'),
    ('6', 'Human Review', 'Recruiter reviews top candidates with full scoring breakdown and call transcripts.'),
]

for col_idx, steps in enumerate([steps_left, steps_right]):
    x = Inches(0.6 + col_idx * 6.0)
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.0 + i * 1.2)
        # Circle
        shape = rect(s, x, y, Inches(0.35), Inches(0.35), fill=RGBColor(0xEE, 0xF2, 0xFF), border=RGBColor(0xEE, 0xF2, 0xFF))
        tb = txtbox(s, x + Inches(0.04), y + Inches(0.04), Inches(0.3), Inches(0.3))
        set_text(tb.text_frame, num, 14, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)
        # Text
        tb = txtbox(s, x + Inches(0.55), y - Inches(0.05), Inches(5.0), Inches(1.0))
        tf = tb.text_frame
        set_text(tf, title, 16, bold=True, color=C['text'])
        add_para(tf, desc, 13, color=C['muted'], space_before=2)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '5/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 6: Architecture ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'ARCHITECTURE', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(8), Inches(0.6))
set_text(tb.text_frame, 'Clean, scalable multi-agent design', 30, bold=True, color=C['text'])

# Diagram area - simple text layout of the architecture
diagram_lines = [
    '+----------------------------------+',
    '|       Recruiter Dashboard        |',
    '+---------------+------------------+',
    '                |',
    '+---------------v------------------+',
    '|         Orchestrator Agent        |',
    '| (Task Router + State + Escalation)|',
    '+--+--------+--------+--------+----+',
    '   |        |        |        |',
    '   v        v        v        v',
    '+--+---+  +-+-----+ ++-----+ +----+',
    '|Sourcer|  |Matcher| |Screen| |Sched|',
    '+-------+  +-------+ +-----+ +-----+',
    '   |        |        |        |',
    '   +--------+--------+--------+',
    '                |',
    '+---------------v-------------------+',
    '|   Memory Layer (Qdrant + Redis    |',
    '|          + PostgreSQL)             |',
    '+-----------------------------------+',
]
tb = txtbox(s, Inches(0.6), Inches(1.8), Inches(6.5), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = False
set_text(tf, '\n'.join(diagram_lines), 11, color=C['muted'])

# Tech stack boxes on right
stacks = [('Framework', 'CrewAI'), ('LLM', 'Gemini 1.5 Pro'), ('Voice', 'Sarvam AI'),
          ('Vector DB', 'Qdrant'), ('Embeddings', 'BGE-M3'), ('Backend', 'FastAPI')]
for i, (label, val) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = Inches(7.5 + col * 2.4)
    y = Inches(1.8 + row * 1.1)
    rect(s, x, y, Inches(2.2), Inches(0.8), fill=C['card_bg'], border=C['card_border'])
    tb = txtbox(s, x + Inches(0.12), y + Inches(0.08), Inches(2.0), Inches(0.65))
    tf = tb.text_frame
    set_text(tf, label, 11, bold=True, color=C['text'])
    add_para(tf, val, 12, color=C['muted'])

tb = txtbox(s, Inches(7.5), Inches(5.6), Inches(5), Inches(0.6))
set_text(tb.text_frame, 'Each agent is independently deployable. Async messaging via Redis pub/sub. Horizontal scaling on every layer.', 11, color=C['light'])

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '6/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 7: Key Differentiator ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xF0, 0xFD, 0xF4))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'KEY DIFFERENTIATOR', 10, bold=True, color=C['green'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.7))
tf = tb.text_frame
set_text(tf, "Voice screening that speaks ", 30, bold=True, color=C['text'])
add_para(tf, "India's language", 30, bold=True, color=C['green'])

# Left card: traditional
rect(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(3.0), fill=C['card_bg'], border=C['card_border'])
tb = txtbox(s, Inches(0.9), Inches(2.15), Inches(5.0), Inches(2.7))
tf = tb.text_frame
set_text(tf, 'Traditional screening', 16, bold=True, color=C['muted'])
for line in [
    '  Text-only, English-only, rigid Q&A',
    '  Filters out great candidates who lack English fluency',
    '  Misses communication nuance, confidence, personality',
    '  Feels robotic and impersonal',
]:
    add_para(tf, line, 13, color=C['muted'], space_before=6)

# Right card: our approach
rect(s, Inches(6.5), Inches(2.0), Inches(5.5), Inches(3.0), fill=C['card_bg'], border=RGBColor(0xBB, 0xF7, 0xD0))
tb = txtbox(s, Inches(6.8), Inches(2.15), Inches(5.0), Inches(2.7))
tf = tb.text_frame
set_text(tf, 'Our Screener Agent', 16, bold=True, color=C['green'])
for line in [
    '  Natural conversation in Hinglish, Hindi, English',
    '  Dynamic questioning that adapts to responses',
    '  Evaluates communication, confidence, cultural fit',
    '  Removes English fluency bias from early filtering',
]:
    add_para(tf, line, 13, color=C['muted'], space_before=6)

# Bottom metrics
for i, (num, label) in enumerate([('80%', 'of Indians prefer\nnative language'), ('+34%', 'accuracy with voice\nvs text-only'), ('2x', 'candidate pool from\nlanguage inclusion')]):
    x = Inches(0.6 + i * 4.2)
    y = Inches(5.3)
    tb = txtbox(s, x, y, Inches(3.5), Inches(1.2))
    tf = tb.text_frame
    set_text(tf, num, 28, bold=True, color=C['green'])
    add_para(tf, label, 11, color=C['light'], space_before=2)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '7/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 8: Impact ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'IMPACT', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(8), Inches(0.6))
set_text(tb.text_frame, 'What this means for a hiring team', 30, bold=True, color=C['text'])

impacts = [
    ('Time to shortlist', '5\u20137 days', '10 minutes'),
    ('Recruiter screening time', '60% of day', '10% of day'),
    ('Candidates evaluated per role', '50\u2013100', '500+'),
    ('Candidate experience', 'Generic, English-only', 'Personalized, multilingual'),
    ('Decision transparency', 'Black box', 'Full reasoning per candidate'),
]
y = Inches(2.0)
for label, before, after in impacts:
    tb = txtbox(s, Inches(0.6), y, Inches(4.5), Inches(0.5))
    set_text(tb.text_frame, label, 14, color=C['text'])
    tb = txtbox(s, Inches(5.5), y, Inches(2.0), Inches(0.5))
    set_text(tb.text_frame, before, 13, color=C['light'], align=PP_ALIGN.RIGHT)
    tb = txtbox(s, Inches(7.8), y, Inches(0.3), Inches(0.5))
    set_text(tb.text_frame, '\u2192', 16, color=C['card_border'], align=PP_ALIGN.CENTER)
    tb = txtbox(s, Inches(8.3), y, Inches(2.5), Inches(0.5))
    set_text(tb.text_frame, after, 14, bold=True, color=C['green'])
    y += Inches(0.55)

# Bottom metric
y = Inches(5.2)
rect(s, Inches(0.6), y, Inches(11.5), Inches(0.7), fill=C['card_bg'], border=C['card_border'])
tb = txtbox(s, Inches(0.9), y + Inches(0.08), Inches(11), Inches(0.55))
tf = tb.text_frame
set_text(tf, '~70% reduction in cost-per-hire', 20, bold=True, color=C['text'], align=PP_ALIGN.CENTER)
add_para(tf, 'for mid-to-senior roles', 11, color=C['light'], align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '8/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 9: Market ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xF0, 0xFD, 0xF4))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'MARKET OPPORTUNITY', 10, bold=True, color=C['green'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.6))
tf = tb.text_frame
set_text(tf, 'Built for Redrob, relevant for ', 30, bold=True, color=C['text'])
add_para(tf, 'everyone', 30, bold=True, color=C['green'])

go_to_market = [
    ('Redrob Integration', 'Native AI hiring layer for India\'s AI OS. Directly aligns with Redrob\'s mission as India\'s AI super app for work.'),
    ('Enterprise SaaS', 'Standalone product for Indian companies (200\u20135000 employees). Plugs into existing ATS systems.'),
    ('API-First', 'Embeddable agent swarm for any ATS. White-label options for recruitment agencies and platforms.'),
]
for i, (title, desc) in enumerate(go_to_market):
    x = Inches(0.6 + i * 4.1)
    y = Inches(2.0)
    rect(s, x, y, Inches(3.8), Inches(2.2), fill=C['card_bg'], border=C['card_border'])
    tb = txtbox(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(1.9))
    tf = tb.text_frame
    set_text(tf, title, 16, bold=True, color=C['text'])
    add_para(tf, desc, 13, color=C['muted'], space_before=8)

# Bottom metrics
for i, (num, label) in enumerate([('$2B+', 'India recruitment\nsoftware TAM by 2027'), ('500', 'Mid-size Indian\ncompanies — initial target')]):
    x = Inches(2.0 + i * 5.5)
    y = Inches(4.6)
    tb = txtbox(s, x, y, Inches(4), Inches(1.0))
    tf = tb.text_frame
    set_text(tf, num, 28, bold=True, color=C['accent'] if i == 0 else C['green'], align=PP_ALIGN.CENTER)
    add_para(tf, label, 11, color=C['light'], align=PP_ALIGN.CENTER, space_before=2)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '9/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 10: Roadmap ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xEE, 0xF2, 0xFF))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'ROADMAP', 10, bold=True, color=C['accent'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(8), Inches(0.6))
set_text(tb.text_frame, 'Beyond the hackathon', 30, bold=True, color=C['text'])

roadmap = [
    ('Q3 2026', 'PoC', 'Core swarm with synthetic data.\nOrchestrator + Matcher + mock sources.', C['accent']),
    ('Q4 2026', 'MVP', 'Live with 5 beta customers.\nLinkedIn source, 3 agents.', C['green']),
    ('Q1 2027', 'Growth', '4+ sources, all 5 agents.\nHinglish voice screening, production.', C['amber']),
    ('Q2 2027', 'Scale', 'Multi-lingual voice, learning loop.\nEnterprise features, API marketplace.', C['pink']),
]
for i, (quarter, title, desc, clr) in enumerate(roadmap):
    x = Inches(0.6 + i * 3.1)
    y = Inches(2.0)
    rect(s, x, y, Inches(2.9), Inches(3.2), fill=C['card_bg'], border=C['card_border'])
    # colored top border
    rect(s, x, y, Inches(2.9), Inches(0.06), fill=clr)
    tb = txtbox(s, x + Inches(0.18), y + Inches(0.2), Inches(2.6), Inches(2.8))
    tf = tb.text_frame
    set_text(tf, quarter, 11, bold=True, color=clr)
    add_para(tf, title, 16, bold=True, color=C['text'], space_before=4)
    add_para(tf, desc, 12, color=C['muted'], space_before=8)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '10/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 11: About Me ----------
s = add_slide(prs)

rect(s, Inches(0.6), Inches(0.5), Inches(1.5), Inches(0.32), fill=RGBColor(0xF0, 0xFD, 0xF4))
tb = txtbox(s, Inches(0.7), Inches(0.52), Inches(1.3), Inches(0.3))
set_text(tb.text_frame, 'ABOUT ME', 10, bold=True, color=C['green'])

tb = txtbox(s, Inches(0.6), Inches(1.1), Inches(10), Inches(0.6))
tf = tb.text_frame
set_text(tf, 'Builder, not just a ', 30, bold=True, color=C['text'])
add_para(tf, 'presenter', 30, bold=True, color=C['green'])

# Profile card
rect(s, Inches(0.6), Inches(2.0), Inches(11.5), Inches(2.0), fill=C['card_bg'], border=C['card_border'])
# Avatar circle
shape = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(2.3), Inches(1.0), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
shape.line.fill.background()
tb = txtbox(s, Inches(1.0), Inches(2.5), Inches(1.0), Inches(0.6))
set_text(tb.text_frame, 'PB', 24, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(2.4), Inches(2.2), Inches(9), Inches(1.6))
tf = tb.text_frame
set_text(tf, 'Class 12 graduate \u00b7 Proficient with AI tools', 16, bold=True, color=C['text'])
add_para(tf, "I just experienced the hiring process as a candidate. I know firsthand how broken it feels on the other side. This isn't theoretical \u2014 it's personal.", 14, color=C['muted'], space_before=6)
add_para(tf, 'Building for India \u2014 where I live, study, and will work.', 14, bold=True, color=C['accent'], space_before=8)

# Quote
y = Inches(4.4)
rect(s, Inches(2.0), y, Inches(8.5), Inches(0.7), fill=C['card_bg'], border=C['card_border'])
tb = txtbox(s, Inches(2.3), y + Inches(0.08), Inches(8), Inches(0.55))
set_text(tb.text_frame, "\u201cI'm not here to pitch. I'm here to build what next India runs on.\u201d", 16, bold=True, color=C['text'], align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(11.8), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '11/12', 11, color=C['light'], align=PP_ALIGN.RIGHT)

# ---------- SLIDE 12: Thank You ----------
s = add_slide(prs)

tb = txtbox(s, Inches(1), Inches(2.0), Inches(10), Inches(1.5))
tf = tb.text_frame
set_text(tf, 'Thank ', 48, bold=True, color=C['text'], align=PP_ALIGN.CENTER)
p = add_para(tf, 'You', 48, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(1), Inches(3.5), Inches(10), Inches(0.5))
set_text(tb.text_frame, "Let's build India's AI-native hiring layer", 20, color=C['muted'], align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(1), Inches(4.5), Inches(10), Inches(1.2))
tf = tb.text_frame
set_text(tf, 'GitHub:  github.com/peter14l/india-runs-ai-recruiter', 14, color=C['text'], align=PP_ALIGN.CENTER)
add_para(tf, 'Event:  India Runs by Redrob AI', 14, color=C['muted'], align=PP_ALIGN.CENTER, space_before=4)
add_para(tf, 'Track:  Challenge 1 \u2014 AI Systems Architect', 14, color=C['muted'], align=PP_ALIGN.CENTER, space_before=4)

# Tags
for i, (tag, clr) in enumerate([('#IndiaRuns', C['accent']), ('#RedrobAI', C['green']), ('#AIAgentSwarm', C['accent'])]):
    x = Inches(4.0 + i * 2.0)
    y = Inches(6.0)
    rect(s, x, y, Inches(1.6), Inches(0.35), fill=RGBColor(0xEE, 0xF2, 0xFF) if i != 1 else RGBColor(0xF0, 0xFD, 0xF4))
    tb = txtbox(s, x + Inches(0.05), y + Inches(0.03), Inches(1.5), Inches(0.3))
    set_text(tb.text_frame, tag, 11, bold=True, color=clr, align=PP_ALIGN.CENTER)

tb = txtbox(s, Inches(0.6), Inches(6.8), Inches(1), Inches(0.4))
set_text(tb.text_frame, '12/12', 11, color=C['light'])

# Save
prs.save(r'E:\india-runs-ai-recruiter\assets\AI_Recruiter_Agent_Swarm_Pitch_Deck.pptx')
print('PPTX saved: E:\\india-runs-ai-recruiter\\assets\\AI_Recruiter_Agent_Swarm_Pitch_Deck.pptx')
